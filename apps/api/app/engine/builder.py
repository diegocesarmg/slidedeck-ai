"""
PPTX Builder — converts the IR (Presentation model) into an actual .pptx file.

Supports three modes:
- From Scratch: Creates a blank presentation
- Template: Opens an existing .pptx and populates it with new content
- Reference: Creates from scratch but applies extracted design tokens
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Optional
from urllib.request import urlopen

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from app.models.slide import (
    Presentation,
    Slide,
    TextBox,
    ImageElement,
    ChartElement,
    HorizontalAlignment,
    VerticalAlignment,
)

logger = logging.getLogger(__name__)

# ── Alignment mappings ───────────────────────────────────────────────────────

_H_ALIGN_MAP = {
    HorizontalAlignment.LEFT: PP_ALIGN.LEFT,
    HorizontalAlignment.CENTER: PP_ALIGN.CENTER,
    HorizontalAlignment.RIGHT: PP_ALIGN.RIGHT,
}

_V_ALIGN_MAP = {
    VerticalAlignment.TOP: MSO_ANCHOR.TOP,
    VerticalAlignment.MIDDLE: MSO_ANCHOR.MIDDLE,
    VerticalAlignment.BOTTOM: MSO_ANCHOR.BOTTOM,
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex color string like '#1a73e8' to an RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


class PresentationBuilder:
    """Builds a python-pptx Presentation from the IR model."""

    # Standard 16:9 widescreen dimensions
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(self) -> None:
        self._pptx: Optional[PptxPresentation] = None

    def build(
        self,
        presentation: Presentation,
        output_path: Path,
        template_path: Path | None = None,
    ) -> Path:
        """
        Build a .pptx file from the Presentation IR model.

        Args:
            presentation: The IR model describing the presentation.
            output_path: Where to save the generated .pptx file.
            template_path: Optional path to a template .pptx. If provided,
                           the template's slide masters are preserved and
                           new slides are added using the template's layouts.

        Returns:
            The path to the saved .pptx file.
        """
        if template_path and template_path.exists():
            logger.info("Building from template: %s", template_path)
            self._pptx = PptxPresentation(str(template_path))
            # Remove existing slides from template (keep masters only)
            self._clear_slides()
        else:
            self._pptx = PptxPresentation()
            self._pptx.slide_width = self.SLIDE_WIDTH
            self._pptx.slide_height = self.SLIDE_HEIGHT

        for slide_model in presentation.slides:
            self._add_slide(slide_model)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        self._pptx.save(str(output_path))
        logger.info("Saved presentation to %s", output_path)
        return output_path

    # ── Private helpers ──────────────────────────────────────────────────

    def _clear_slides(self) -> None:
        """Remove all existing slides from the presentation (preserves masters)."""
        slide_ids = list(self._pptx.slides._sldIdLst)
        for sld_id in slide_ids:
            rId = sld_id.get("r:id") or sld_id.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            if rId:
                self._pptx.part.drop_rel(rId)
            self._pptx.slides._sldIdLst.remove(sld_id)

    def _get_layout(self, slide_model: Slide):
        """Get the best matching layout from available slide layouts."""
        # Try to find a layout matching the model's layout name
        layout_name_map = {
            "title": ["Title Slide", "Title"],
            "title_content": ["Title and Content", "Title, Content"],
            "two_column": ["Two Content", "Comparison"],
            "section_header": ["Section Header"],
            "blank": ["Blank"],
            "image_full": ["Blank", "Picture with Caption"],
        }

        target_names = layout_name_map.get(slide_model.layout.value, ["Blank"])

        for layout in self._pptx.slide_layouts:
            if layout.name in target_names:
                return layout

        # Fallback: use the last layout (typically Blank) or index 6
        try:
            return self._pptx.slide_layouts[6]
        except IndexError:
            return self._pptx.slide_layouts[-1]

    def _add_slide(self, slide_model: Slide) -> None:
        """Add a single slide to the presentation."""
        slide_layout = self._get_layout(slide_model)
        slide = self._pptx.slides.add_slide(slide_layout)

        # Set background color
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(slide_model.background_color)

        # Add each element
        for element in slide_model.elements:
            if isinstance(element, TextBox):
                self._add_textbox(slide, element)
            elif isinstance(element, ImageElement):
                self._add_image(slide, element)
            elif isinstance(element, ChartElement):
                self._add_chart_placeholder(slide, element)

        # Speaker notes
        if slide_model.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_model.speaker_notes

    def _add_textbox(self, slide, tb: TextBox) -> None:
        """Add a text box to the slide."""
        txBox = slide.shapes.add_textbox(
            Inches(tb.x), Inches(tb.y), Inches(tb.width), Inches(tb.height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Vertical alignment
        tf.auto_size = None
        if tb.vertical_alignment in _V_ALIGN_MAP:
            txBox.text_frame._txBody.bodyPr.set(
                "anchor", tb.vertical_alignment.value
            )

        p = tf.paragraphs[0]
        p.text = tb.content

        # Paragraph alignment
        p.alignment = _H_ALIGN_MAP.get(tb.alignment, PP_ALIGN.LEFT)

        # Font settings
        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs:
            run.text = tb.content
        run.font.name = tb.font_name
        run.font.size = Pt(tb.font_size)
        run.font.bold = tb.font_bold
        run.font.italic = tb.font_italic
        run.font.color.rgb = _hex_to_rgb(tb.font_color)

    def _add_image(self, slide, img: ImageElement) -> None:
        """Add an image to the slide (from URL or local path)."""
        try:
            if img.url:
                # Download from URL
                image_data = urlopen(img.url).read()
                image_stream = io.BytesIO(image_data)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(img.x),
                    Inches(img.y),
                    Inches(img.width),
                    Inches(img.height),
                )
            elif img.path:
                slide.shapes.add_picture(
                    img.path,
                    Inches(img.x),
                    Inches(img.y),
                    Inches(img.width),
                    Inches(img.height),
                )
            else:
                logger.warning("Image element has neither url nor path, skipping.")
        except Exception as e:
            logger.error("Failed to add image: %s", e)
            # Add a placeholder textbox instead
            self._add_textbox(
                slide,
                TextBox(
                    content=f"[Image: {img.alt_text or 'unavailable'}]",
                    x=img.x,
                    y=img.y,
                    width=img.width,
                    height=img.height,
                    font_size=14,
                    font_color="#999999",
                    alignment=HorizontalAlignment.CENTER,
                    vertical_alignment=VerticalAlignment.MIDDLE,
                ),
            )

    def _add_chart_placeholder(self, slide, chart: ChartElement) -> None:
        """
        Add a chart as a styled table (python-pptx chart support is limited).

        Renders chart data as a table with a title row.
        """
        if not chart.categories or not chart.series:
            logger.warning("Chart has no data, adding title only.")
            self._add_textbox(
                slide,
                TextBox(
                    content=chart.title or "[Chart]",
                    x=chart.x,
                    y=chart.y,
                    width=chart.width,
                    height=0.5,
                    font_size=16,
                    font_bold=True,
                    alignment=HorizontalAlignment.CENTER,
                ),
            )
            return

        rows = len(chart.series) + 1  # header + data rows
        cols = len(chart.categories) + 1  # label col + data cols

        table_shape = slide.shapes.add_table(
            rows, cols, Inches(chart.x), Inches(chart.y),
            Inches(chart.width), Inches(chart.height),
        )
        table = table_shape.table

        # Header row: empty cell + category names
        table.cell(0, 0).text = chart.title or ""
        for ci, cat in enumerate(chart.categories):
            table.cell(0, ci + 1).text = str(cat)

        # Data rows
        for ri, s in enumerate(chart.series):
            table.cell(ri + 1, 0).text = s.get("name", "")
            for ci, val in enumerate(s.get("values", [])):
                table.cell(ri + 1, ci + 1).text = str(val)
