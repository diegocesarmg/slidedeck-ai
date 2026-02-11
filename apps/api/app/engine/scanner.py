"""
Scanner — extracts design tokens from an existing .pptx file.

Used by Template Mode and Reference Mode to analyze visual style
and inject it as constraints into the LLM prompt.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from app.models.slide import DesignTokens

logger = logging.getLogger(__name__)


def _rgb_to_hex(rgb: RGBColor | None) -> str | None:
    """Convert an RGBColor to a hex string like '#1a73e8'."""
    if rgb is None:
        return None
    return f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"


def _extract_colors(pptx: PptxPresentation) -> list[str]:
    """Extract unique colors used across all slides."""
    colors: set[str] = set()

    for slide in pptx.slides:
        # Background color
        bg = slide.background
        if bg.fill and bg.fill.type is not None:
            try:
                hex_color = _rgb_to_hex(bg.fill.fore_color.rgb)
                if hex_color:
                    colors.add(hex_color)
            except (AttributeError, TypeError):
                pass

        # Shape colors
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        try:
                            if run.font.color and run.font.color.rgb:
                                hex_color = _rgb_to_hex(run.font.color.rgb)
                                if hex_color:
                                    colors.add(hex_color)
                        except (AttributeError, TypeError):
                            pass

            # Shape fill
            if hasattr(shape, "fill") and shape.fill.type is not None:
                try:
                    hex_color = _rgb_to_hex(shape.fill.fore_color.rgb)
                    if hex_color:
                        colors.add(hex_color)
                except (AttributeError, TypeError):
                    pass

    return sorted(colors)


def _extract_fonts(pptx: PptxPresentation) -> list[str]:
    """Extract unique font names used across all slides."""
    fonts: set[str] = set()

    for slide in pptx.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.add(run.font.name)

    return sorted(fonts)


def _extract_layouts(pptx: PptxPresentation) -> list[str]:
    """Extract available slide layout names from the template."""
    layouts: list[str] = []
    for layout in pptx.slide_layouts:
        if layout.name:
            layouts.append(layout.name)
    return layouts


def scan_pptx(file_path: Path) -> DesignTokens:
    """
    Analyze a .pptx file and extract its design tokens.

    Args:
        file_path: Path to the .pptx file to analyze.

    Returns:
        A DesignTokens model with extracted colors, fonts, and layouts.
    """
    logger.info("Scanning .pptx: %s", file_path)
    pptx = PptxPresentation(str(file_path))

    colors = _extract_colors(pptx)
    fonts = _extract_fonts(pptx)
    layouts = _extract_layouts(pptx)

    # Determine primary/secondary colors (heuristic: first two most-used)
    primary_color = colors[0] if len(colors) > 0 else "#1a73e8"
    secondary_color = colors[1] if len(colors) > 1 else "#e8710a"
    background_color = "#FFFFFF"  # default, hard to detect reliably

    # Determine heading/body fonts (heuristic: first two)
    font_heading = fonts[0] if len(fonts) > 0 else "Calibri"
    font_body = fonts[1] if len(fonts) > 1 else fonts[0] if fonts else "Calibri"

    tokens = DesignTokens(
        primary_color=primary_color,
        secondary_color=secondary_color,
        background_color=background_color,
        font_heading=font_heading,
        font_body=font_body,
        layout_names=layouts,
        extracted_colors=colors,
        extracted_fonts=fonts,
    )

    logger.info(
        "Extracted tokens — colors: %d, fonts: %d, layouts: %d",
        len(colors),
        len(fonts),
        len(layouts),
    )
    return tokens
